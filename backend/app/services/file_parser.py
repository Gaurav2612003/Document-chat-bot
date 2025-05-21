from pdfminer.high_level import extract_text
import os
from pdfminer.high_level import extract_text
from docx import Document as DocxDocument
from pptx import Presentation
from pdfminer.high_level import extract_pages
from pdfminer.layout import LTTextContainer



def extract_text_from_pdf(file_path: str):
    page_data = []
    for page_num, page_layout in enumerate(extract_pages(file_path), start=1):
        page_text = ""
        for element in page_layout:
            if isinstance(element, LTTextContainer):
                page_text += element.get_text()
        if page_text.strip():
            page_data.append({
                "text": page_text.strip(),
                "metadata": {
                    "filename": os.path.basename(file_path),
                    "source_type": "pdf",
                    "page_number": page_num
                }
            })
    return page_data


def extract_text_from_docx(file_path: str) -> str:
    doc = DocxDocument(file_path)
    return "\n".join([para.text for para in doc.paragraphs if para.text.strip()])

def extract_text_from_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return "\n".join(text_runs)

def extract_text_from_txt(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8") as f:
        return f.read()

def extract_paragraphs_with_metadata(content: str, source_type: str, filename: str):
    paragraphs = [p.strip() for p in content.split('\n') if p.strip()]
    metadata_paragraphs = []
    for idx, paragraph in enumerate(paragraphs):
        metadata_paragraphs.append({
            "text": paragraph,
            "metadata": {
                "filename": filename,
                "source_type": source_type,
                "paragraph_number": idx + 1
            }
        })
    return metadata_paragraphs